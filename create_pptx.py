#!/usr/bin/env python3
"""Generate EcommSpringBot project briefing presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
DARK_BG    = RGBColor(0x06, 0x5A, 0x82)   # deep blue
TEAL       = RGBColor(0x1C, 0x72, 0x93)   # teal
ACCENT     = RGBColor(0x4E, 0xCD, 0xC4)   # mint accent
LIGHT_BG   = RGBColor(0xF0, 0xF8, 0xFB)   # very light blue bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x1A, 0x23, 0x32)   # dark text
TEXT_MUTED  = RGBColor(0x64, 0x74, 0x8B)   # muted text
CARD_BG    = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT_LIGHT = RGBColor(0xCC, 0xF3, 0xF0)  # very light mint
DARK_BG2    = RGBColor(0x04, 0x3C, 0x56)   # slightly darker for variety

# ── Presentation setup ─────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ───────────────────────────────────────────

def add_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text,
                 font_name='Calibri', font_size=Pt(14), color=TEXT_DARK,
                 bold=False, alignment=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                 line_spacing=1.15):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size.pt * line_spacing)
    tf.paragraphs[0].space_before = Pt(0)
    return txBox

def add_card(slide, left, top, width, height, title, body,
             title_color=TEAL, body_color=TEXT_MUTED, bg_color=WHITE,
             shadow=True):
    """Add a card with rounded rectangle bg, title, and body text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    # rounded corners
    shape.adjustments[0] = 0.05

    # Title
    add_text_box(slide, left + 0.25, top + 0.2, width - 0.5, 0.45,
                 title, font_name='Cambria', font_size=Pt(16),
                 color=title_color, bold=True)
    # Body
    add_text_box(slide, left + 0.25, top + 0.65, width - 0.5, height - 0.85,
                 body, font_size=Pt(12), color=body_color)
    return shape

def add_slide_title(slide, title, subtitle=None):
    """Add consistent slide title with optional subtitle."""
    add_text_box(slide, 0.8, 0.45, 11.7, 0.65, title,
                 font_name='Cambria', font_size=Pt(32), color=TEXT_DARK,
                 bold=True)
    if subtitle:
        add_text_box(slide, 0.8, 1.05, 11.7, 0.35, subtitle,
                     font_name='Calibri', font_size=Pt(14), color=TEXT_MUTED)

def add_dot(slide, left, top, size, color):
    """Add a small colored circle as a decorative dot/accent."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, color, thickness=Pt(1.5)):
    """Add a horizontal line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Pt(thickness)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_arrow_right(slide, left, top, width, color):
    """Add a right-pointing chevron."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(left), Inches(top),
        Inches(width), Inches(0.25)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_BG)

# Large decorative circle - top right
circle1 = slide1.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5)
)
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(0x07, 0x6E, 0x99)
circle1.line.fill.background()

# Small accent circle
circle2 = slide1.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(1.2), Inches(5.2), Inches(1.2), Inches(1.2)
)
circle2.fill.solid()
circle2.fill.fore_color.rgb = ACCENT
circle2.line.fill.background()
# set transparency manually
circle2.fill.fore_color.brightness = 0.0

# Main title
add_text_box(slide1, 1.2, 1.8, 11, 1.2,
             "EcommSpringBot",
             font_name='Cambria', font_size=Pt(54), color=WHITE, bold=True)

# Subtle line under title
add_line(slide1, 1.2, 3.05, 2.5, ACCENT, Pt(3))

# Subtitle
add_text_box(slide1, 1.2, 3.35, 10, 0.55,
             "基于 Spring AI Alibaba 的智能电商客服助手系统",
             font_name='Calibri', font_size=Pt(22), color=RGBColor(0xCC, 0xEC, 0xF2))

# Description
add_text_box(slide1, 1.2, 4.1, 9, 0.8,
             "集订单管理、知识问答、多层记忆、人机协同于一体的\n企业级电商智能助手平台",
             font_name='Calibri', font_size=Pt(16), color=RGBColor(0x99, 0xD6, 0xE6))

# Date
add_text_box(slide1, 1.2, 6.0, 4, 0.4,
             "2026年7月  ·  项目汇报",
             font_name='Calibri', font_size=Pt(13), color=RGBColor(0x88, 0xC4, 0xD6))


# ══════════════════════════════════════════════════════════════
# SLIDE 2: 项目概述
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide2, LIGHT_BG)

add_slide_title(slide2, "项目概述",
                "EcommSpringBot 是什么？能做什么？")

# 4 cards in a 2x2 grid
cards_data = [
    ("订单智能管理", "支持订单查询、取消、创建等\n操作，通过自然语言即可完成\n电商订单全生命周期管理"),
    ("知识库智能问答", "内置 7 大领域知识库，涵盖\nHR、财务、电商规则、客服等\n场景，RAG 检索精准回答"),
    ("多层记忆系统", "短期记忆（Redis）+ 长期记忆\n（Milvus），自动提取和整合\n用户画像、偏好与历史事实"),
    ("人机协同审批", "敏感操作（退款、取消）触发\n人工审批节点，StateGraph\n支持中断-恢复的人机交互"),
]

card_w = 5.3
card_h = 2.0
start_x = 1.2
start_y = 1.75
gap_x = 0.5
gap_y = 0.35

for i, (title, body) in enumerate(cards_data):
    col = i % 2
    row = i // 2
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_card(slide2, x, y, card_w, card_h, title, body)

# Bottom accent dots
for j in range(4):
    add_dot(slide2, 1.2 + j * 0.5, 6.8, 0.12, ACCENT if j == 0 else CARD_BORDER)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: 技术架构
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

add_slide_title(slide3, "技术架构",
                "核心技术栈一览")

# Tech stack sections - 3 columns
col_data = [
    {
        "title": "AI & 模型层",
        "color": TEAL,
        "items": [
            "Spring AI Alibaba",
            "DashScope (Qwen)",
            "text-embedding-v2",
            "qwen3-rerank",
        ]
    },
    {
        "title": "基础设施层",
        "color": ACCENT,
        "items": [
            "MySQL (业务数据)",
            "Redis (短期记忆)",
            "Milvus (向量存储)",
            "RocketMQ (消息)",
            "Elasticsearch (追踪)",
        ]
    },
    {
        "title": "框架与协议",
        "color": RGBColor(0xF9, 0xA8, 0x26),
        "items": [
            "Spring Boot 3.5",
            "MCP 协议",
            "StateGraph 编排",
            "ReAct Agent",
            "WebFlux / SSE",
        ]
    },
]

col_w = 3.5
col_start_x = 1.2
col_gap = 0.55
col_y = 1.7

for i, col in enumerate(col_data):
    x = col_start_x + i * (col_w + col_gap)

    # Column header
    header_shape = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(col_y), Inches(col_w), Inches(0.55)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = col["color"]
    header_shape.line.fill.background()
    header_shape.adjustments[0] = 0.15

    add_text_box(slide3, x, col_y + 0.06, col_w, 0.45,
                 col["title"], font_name='Cambria', font_size=Pt(15),
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)

    # Items
    item_y = col_y + 0.8
    for j, item in enumerate(col["items"]):
        # Dot indicator
        add_dot(slide3, x + 0.3, item_y + j * 0.38 + 0.14, 0.1, col["color"])
        add_text_box(slide3, x + 0.6, item_y + j * 0.38, col_w - 0.8, 0.35,
                     item, font_size=Pt(13), color=TEXT_DARK)

# Bottom section - 一行列出关键版本号
add_line(slide3, 1.2, 6.0, 11.15, CARD_BORDER)
add_text_box(slide3, 1.2, 6.15, 11, 0.35,
             "Java 17  ·  Spring Boot 3.5.7  ·  Spring AI 1.1.0  ·  Milvus 2.5  ·  Qwen (Plus / Max / Turbo)",
             font_name='Calibri', font_size=Pt(11), color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: 核心模块
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, LIGHT_BG)

add_slide_title(slide4, "核心模块架构",
                "8 大模块组成微服务架构体系")

# 8 modules displayed as connected blocks
# We'll use a central flow layout
modules = [
    ("mall-order", "订单核心服务\nMySQL + MyBatis\nCRUD 订单管理", "8081"),
    ("mall-order-\ncmp-server", "MCP 服务中间层\nWebFlux SSE\n工具暴露与代理", "8082"),
    ("mall-order-\ncmp-sse-client", "Agent 接入层\nReAct Agent\nSSE + Redis 记忆", "8888"),
    ("mall-order-\nmilvus-rag", "Milvus RAG\n向量检索 + 重排序\n7 领域知识库", "8086"),
    ("mall-order-\nmemory", "多层记忆库\nRedis + Milvus\n自动抽取整合", "—"),
    ("mall-order-\nagent", "订单 Agent\nStateGraph 编排\n人机协同审批", "8087"),
    ("mall-order-\nobservability", "可观测性\nRocketMQ 传输\nES 追踪存储", "—"),
    ("skill-agent-\ndemo", "Skill Agent\n技能系统演示\nArxiv / Web 搜索", "—"),
]

# Layout: 4 rows x 2 columns
mod_w = 5.3
mod_h = 1.1
mod_start_x = 1.2
mod_start_y = 1.6
mod_gap_x = 0.6
mod_gap_y = 0.3

for i, (name, desc, port) in enumerate(modules):
    col = i % 2
    row = i // 2
    x = mod_start_x + col * (mod_w + mod_gap_x)
    y = mod_start_y + row * (mod_h + mod_gap_y)

    # Module card background
    card = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(mod_w), Inches(mod_h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(0.5)
    card.adjustments[0] = 0.06

    # Number badge on left
    badge = slide4.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.2), Inches(y + 0.3), Inches(0.5), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.name = 'Cambria'
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Module name
    add_text_box(slide4, x + 0.9, y + 0.1, 2.2, 0.4,
                 name, font_name='Cambria', font_size=Pt(13),
                 color=TEXT_DARK, bold=True)

    # Description
    add_text_box(slide4, x + 0.9, y + 0.5, 3.0, 0.65,
                 desc, font_name='Calibri', font_size=Pt(10),
                 color=TEXT_MUTED, line_spacing=1.3)

    # Port badge
    port_badge = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + mod_w - 1.0), Inches(y + 0.15), Inches(0.8), Inches(0.3)
    )
    port_badge.fill.solid()
    port_badge.fill.fore_color.rgb = ACCENT_LIGHT
    port_badge.line.fill.background()
    port_badge.adjustments[0] = 0.3
    tf2 = port_badge.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f":{port}" if port != "—" else "lib"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(9)
    p2.font.color.rgb = TEAL
    p2.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════
# SLIDE 5: AI Agent 工作流
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)

add_slide_title(slide5, "AI Agent 工作流",
                "基于 StateGraph 的智能编排引擎")

# Flow diagram: nodes connected by arrows
# Top section = legend for 3 execution paths
paths = [
    ("路径一：知识库命中 → 完整 RAG + LLM 生成", TEAL),
    ("路径二：知识库未命中 → 快捷应答", TEXT_MUTED),
    ("路径三：敏感操作 → 人工审批介入", RGBColor(0xF9, 0xA8, 0x26)),
]
for i, (label, color) in enumerate(paths):
    add_dot(slide5, 1.2 + i * 4.0, 1.65, 0.12, color)
    add_text_box(slide5, 1.5 + i * 4.0, 1.6, 3.5, 0.3,
                 label, font_size=Pt(12), color=color)

add_line(slide5, 1.2, 2.05, 11.15, CARD_BORDER)

# Graph nodes - horizontal flow
nodes = [
    ("Planner", "意图解析\n策略规划"),
    ("Action\nRunner", "工具执行\n(RAG/Tool/\nMemory)"),
    ("Prompt\nBuilder", "上下文组装\n提示构建"),
    ("LLM\nNode", "大模型生成\n(qwen-plus)"),
    ("Human\nNode", "人工审批\n中断-恢复"),
    ("Answer\nNode", "结果输出\n记忆持久化"),
]

node_w = 1.65
node_h = 1.0
node_start_x = 0.7
node_y = 2.6
node_gap = 0.45

for i, (name, desc) in enumerate(nodes):
    x = node_start_x + i * (node_w + node_gap)

    # Node rectangle
    color = TEAL if i < 4 else (RGBColor(0xF9, 0xA8, 0x26) if i == 4 else ACCENT)
    node = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(node_y), Inches(node_w), Inches(node_h)
    )
    node.fill.solid()
    node.fill.fore_color.rgb = color
    node.line.fill.background()
    node.adjustments[0] = 0.12

    # Node name
    add_text_box(slide5, x, node_y + 0.2, node_w, 0.35,
                 name, font_name='Cambria', font_size=Pt(12),
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow between nodes (except after last)
    if i < len(nodes) - 1:
        ax = x + node_w + 0.05
        ay = node_y + node_h / 2 - 0.06
        aw = node_gap - 0.15
        arr = slide5.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(ax), Inches(ay), Inches(aw), Inches(0.13)
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = CARD_BORDER
        arr.line.fill.background()

# Description below nodes
node_desc_y = node_y + node_h + 0.35
for i, (name, desc) in enumerate(nodes):
    x = node_start_x + i * (node_w + node_gap)
    add_text_box(slide5, x, node_desc_y, node_w, 1.0,
                 desc, font_size=Pt(11), color=TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.35)

# Bottom note
add_text_box(slide5, 1.2, 6.4, 11, 0.4,
             "完整流程: START → Planner → ActionRunner → PromptBuilder → LLM → [Human] → Answer → END",
             font_size=Pt(11), color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: RAG 知识检索流程
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, LIGHT_BG)

add_slide_title(slide6, "RAG 知识检索流程",
                "两阶段检索 + 重排序，确保回答精准可控")

# Pipeline steps arranged vertically with connecting elements
steps = [
    ("01", "文档摄入", "PDF 上传 / 批量导入\n7 个内置领域知识库\nHR · 财务 · 电商 · 物流 · 客服 · 技术 · AI"),
    ("02", "文本预处理", "PDF 清洗去噪\n章节目录级切割\n250 Token 滑动分块"),
    ("03", "向量化存储", "text-embedding-v2\n1536 维向量\nMilvus 集合 mall_rag_v2"),
    ("04", "检索 + 重排序", "余弦相似度粗筛\nqwen3-rerank 精细排序\n元数据过滤（角色/部门）"),
    ("05", "答案生成", "qwen-plus 生成\n防幻觉严格提示词\n引用来源可追溯"),
]

step_w = 11.15
step_h = 0.85
step_start_x = 1.2
step_start_y = 1.75
step_gap = 0.25

for i, (num, title, desc) in enumerate(steps):
    y = step_start_y + i * (step_h + step_gap)

    # Step background card
    step_card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(step_start_x), Inches(y), Inches(step_w), Inches(step_h)
    )
    step_card.fill.solid()
    step_card.fill.fore_color.rgb = WHITE
    step_card.line.color.rgb = CARD_BORDER
    step_card.line.width = Pt(0.5)
    step_card.adjustments[0] = 0.06

    # Number circle
    num_badge = slide6.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(step_start_x + 0.25), Inches(y + 0.18), Inches(0.46), Inches(0.46)
    )
    num_badge.fill.solid()
    num_badge.fill.fore_color.rgb = TEAL
    num_badge.line.fill.background()
    tf = num_badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.name = 'Cambria'
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Step title
    add_text_box(slide6, step_start_x + 0.95, y + 0.08, 3.0, 0.35,
                 title, font_name='Cambria', font_size=Pt(14),
                 color=TEXT_DARK, bold=True)

    # Step description
    add_text_box(slide6, step_start_x + 0.95, y + 0.38, 9.5, 0.48,
                 desc, font_name='Calibri', font_size=Pt(11),
                 color=TEXT_MUTED, line_spacing=1.3)

    # Connector arrow between steps
    if i < len(steps) - 1:
        connector_y = y + step_h
        arr = slide6.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(step_start_x + step_w / 2 - 0.08), Inches(connector_y - 0.02),
            Inches(0.16), Inches(step_gap + 0.04)
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = CARD_BORDER
        arr.line.fill.background()

# Key metric at bottom
add_text_box(slide6, 1.2, 6.45, 11.15, 0.35,
             "7 个领域  ·  250 Token 分块  ·  1536 维向量  ·  余弦 + 重排序双阶段检索  ·  防幻觉提示词",
             font_size=Pt(11), color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: 多层记忆系统
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, WHITE)

add_slide_title(slide7, "多层记忆系统",
                "短期 + 长期双层记忆，自动抽取与整合")

# Two main columns: short-term and long-term memory
col_w2 = 5.15
col_x1 = 1.2
col_x2 = col_x1 + col_w2 + 0.85
col_top = 2.0

# --- Short-term memory column ---
st_card = slide7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(col_x1), Inches(col_top), Inches(col_w2), Inches(3.8)
)
st_card.fill.solid()
st_card.fill.fore_color.rgb = RGBColor(0xF8, 0xFD, 0xFE)
st_card.line.color.rgb = CARD_BORDER
st_card.line.width = Pt(0.5)
st_card.adjustments[0] = 0.05

# Column header
st_header = slide7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(col_x1 + 0.3), Inches(col_top + 0.3), Inches(4.55), Inches(0.5)
)
st_header.fill.solid()
st_header.fill.fore_color.rgb = TEAL
st_header.line.fill.background()
st_header.adjustments[0] = 0.15
add_text_box(slide7, col_x1 + 0.3, col_top + 0.34, 4.55, 0.45,
             "短期记忆  ·  Redis", font_name='Cambria', font_size=Pt(16),
             color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

st_items = [
    "存储最近的对话上下文",
    "支持 TTL 自动过期策略",
    "默认 7 天保留周期",
    "Redisson 客户端连接池",
    "消息达到阈值自动触发整合",
]
for j, item in enumerate(st_items):
    iy = col_top + 1.15 + j * 0.42
    add_dot(slide7, col_x1 + 0.5, iy + 0.1, 0.1, ACCENT)
    add_text_box(slide7, col_x1 + 0.85, iy, 4.2, 0.35,
                 item, font_size=Pt(12), color=TEXT_DARK)

# --- Long-term memory column ---
lt_card = slide7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(col_x2), Inches(col_top), Inches(col_w2), Inches(3.8)
)
lt_card.fill.solid()
lt_card.fill.fore_color.rgb = RGBColor(0xF8, 0xFD, 0xFE)
lt_card.line.color.rgb = CARD_BORDER
lt_card.line.width = Pt(0.5)
lt_card.adjustments[0] = 0.05

lt_header = slide7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(col_x2 + 0.3), Inches(col_top + 0.3), Inches(4.55), Inches(0.5)
)
lt_header.fill.solid()
lt_header.fill.fore_color.rgb = ACCENT
lt_header.line.fill.background()
lt_header.adjustments[0] = 0.15
add_text_box(slide7, col_x2 + 0.3, col_top + 0.34, 4.55, 0.45,
             "长期记忆  ·  Milvus", font_name='Cambria', font_size=Pt(16),
             color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)

lt_items = [
    "用户画像 (User Profile)",
    "事实记忆 (Facts)",
    "对话摘要 (Summaries)",
    "正则 + LLM 两阶段抽取",
    "定时自动整合 (默认 5 min)",
]
for j, item in enumerate(lt_items):
    iy = col_top + 1.15 + j * 0.42
    add_dot(slide7, col_x2 + 0.5, iy + 0.1, 0.1, TEAL)
    add_text_box(slide7, col_x2 + 0.85, iy, 4.2, 0.35,
                 item, font_size=Pt(12), color=TEXT_DARK)

# Bottom section: consolidation flow
add_line(slide7, 1.2, 6.15, 11.15, CARD_BORDER)
add_text_box(slide7, 1.2, 6.3, 11.15, 0.35,
             "记忆整合: 对话消息 → 阈值触发 (20条) / 定时触发 (5min) → 正则初筛 → LLM 结构化抽取 → Redis + Milvus 双写",
             font_size=Pt(11), color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Central arrow connecting the two
mid_x = col_x1 + col_w2 + 0.425  # center of gap between columns
arrow = slide7.shapes.add_shape(
    MSO_SHAPE.LEFT_RIGHT_ARROW,
    Inches(mid_x), Inches(col_top + 1.6), Inches(0.35), Inches(0.35)
)
arrow.fill.solid()
arrow.fill.fore_color.rgb = TEXT_MUTED
arrow.line.fill.background()

add_text_box(slide7, mid_x - 0.35, col_top + 2.1, 1.05, 0.3,
             "整合", font_size=Pt(10), color=TEXT_MUTED,
             alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: 总结与展望
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, DARK_BG)

# Decorative circle
circle_final = slide8.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(-2), Inches(3.5), Inches(5), Inches(5)
)
circle_final.fill.solid()
circle_final.fill.fore_color.rgb = RGBColor(0x07, 0x6E, 0x99)
circle_final.line.fill.background()

add_text_box(slide8, 1.5, 0.8, 10, 0.8,
             "总结与展望",
             font_name='Cambria', font_size=Pt(40), color=WHITE, bold=True)

add_line(slide8, 1.5, 1.7, 2.5, ACCENT, Pt(3))

# Left side:已完成
add_text_box(slide8, 1.5, 2.2, 5, 0.4,
             "已实现能力", font_name='Cambria', font_size=Pt(20),
             color=ACCENT, bold=True)

done_items = [
    "Spring AI Alibaba + DashScope 多模型集成",
    "MCP 协议工具暴露与 Agent 调用",
    "StateGraph 智能编排 + 人机协同审批",
    "Milvus 向量 RAG + 重排序双阶段检索",
    "Redis + Milvus 多层记忆自动管理",
    "RocketMQ + ES 全链路可观测追踪",
]
for j, item in enumerate(done_items):
    iy = 2.7 + j * 0.45
    add_dot(slide8, 1.5, iy + 0.1, 0.12, ACCENT)
    add_text_box(slide8, 1.85, iy, 5.5, 0.35,
                 item, font_size=Pt(14), color=WHITE)

# Right side: 未来规划
add_text_box(slide8, 7.5, 2.2, 5, 0.4,
             "未来规划", font_name='Cambria', font_size=Pt(20),
             color=RGBColor(0xF9, 0xA8, 0x26), bold=True)

future_items = [
    "多 Agent 协作与对话",
    "Skill 技能市场扩展",
    "多模态支持（图片理解）",
    "A/B 评估框架完善",
    "生产级部署与监控",
    "开放 API 与插件生态",
]
for j, item in enumerate(future_items):
    iy = 2.7 + j * 0.45
    add_dot(slide8, 7.5, iy + 0.1, 0.12, RGBColor(0xF9, 0xA8, 0x26))
    add_text_box(slide8, 7.85, iy, 5.5, 0.35,
                 item, font_size=Pt(14), color=WHITE)

# Bottom line
add_line(slide8, 1.5, 6.5, 10, RGBColor(0x07, 0x6E, 0x99))

add_text_box(slide8, 1.5, 6.65, 10, 0.4,
             "EcommSpringBot  ·  基于 Spring AI Alibaba 的企业级智能电商助手",
             font_name='Calibri', font_size=Pt(13), color=RGBColor(0x88, 0xC4, 0xD6),
             alignment=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────
output_path = "/sessions/stoic-eloquent-wright/mnt/EcommSpringBot/EcommSpringBot_项目汇报.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
